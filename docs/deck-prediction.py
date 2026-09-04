# Дека «Vector Prediction» — тема 01-obsidian-neon (неон-синий на тёмно-синем)
# 12 слайдов: титул, проблема, решение, данные, роли/путь, архитектура, лицензия,
# метрики/бенчмарк, прогноз-выход, каденция, roadmap, финал.
import os
OUT_FMT = os.path.expanduser('~/projects/vector-legal-decks15/vector-prediction-{theme}.pptx')

DATA = dict(
    kicker='ПРОГНОЗНЫЙ ДВИЖОК · OSMOSY VECTOR',
    title='Vector Prediction',
    subtitle='Гибридный пайплайн TimesFM: прод-прогноз на 2.5, исследования на 3.0 — '
             'спрос, продажи, трафик кампаний',
    chips=[('2.5', 'прод · Apache'), ('3.0', 'research'), ('6x', 'быстрее ковариаты')],
    github='github.com/Osmosy/vector-prediction',
    footer_tag='Osmosy · Hermes Agent · 2026',
    intro_lead='Zero-shot прогнозирование спроса, продаж и трафика кампаний без '
               'обучения под задачу — локально на CPU. Два контура, разделённые по лицензии.',
    intro_cards=[
        ('Контур A — ПРОД', ['TimesFM 2.5 · Apache-2.0', 'Боевые прогнозы: планы', 'производства и закупки,', 'KPI кампаний, отчёты']),
        ('Контур B — RESEARCH', ['TimesFM 3.0 · non-commercial', 'Только сравнение моделей:', 'MAE, скорость, ковариаты', 'Не попадает в решения']),
        ('Роли', ['Люди: данные и решения', 'Агенты: сборка, прогноз,', 'ревью качества и лицензии']),
    ],

    # 02 · Проблема
    prob_cards=[
        ('«Как в прошлый раз»', ['План по интуиции и прошлому опыту.', 'Акции и праздники — сюрприз:', 'то дефицит, то перепроизводство']),
        ('Excel вручную', ['Часы на сводные таблицы,', 'средние за месяц — не прогноз.', 'Сезон не учитывается совсем']),
        ('Облачные прогнозы', ['Данные продаж уходят наружу,', 'подписка за каждый прогон,', 'чёрный ящик без интервалов']),
    ],
    prob_panel='Цена ошибки в мясопродуктах',
    prob_lines=[
        'Скоропорт: перепроизводство — прямые списания, недопроизводство — упущенная выручка и пустая полка.',
        '!Шашлычный сезон, Пасха, Новый год — пики, которые нельзя «усреднить»: без сезонной истории модель слепа.',
        'Промо-дни дают +20–40% к продажам — план без календаря акций систематически ошибается именно в эти дни.',
    ],

    # 03 · Решение
    sol_cards=[
        ('Zero-shot', ['Без обучения под задачу:', 'загрузили историю — получили', 'прогноз. 200M параметров,', 'претрейн на млрд точек']),
        ('Интервалы, не числа', ['Каждый день: cautious /', 'forecast / optimistic —', 'скоропорт и заморозка', 'планируются по-разному']),
        ('Локально', ['Веса на своём железе:', 'данные продаж не покидают', 'контур. CPU 20 ядер,', 'RAM ~1.5 ГБ']),
    ],
    sol_panel='Главный эффект',
    sol_lines=[
        '!Модель «видит» будущую акцию: календарь промо подаётся как past-future ковариат — рост заложен заранее.',
        'Кросс-связи рядов: промо в сосисках видно и в прогнозе ветчины (мультивариантность 3.0; в 2.5 — через XReg).',
        'Аномалии в истории ловятся квантилями: факт вне 80% интервала = WARNING/CRITICAL сигнал.',
    ],

    # 04 · Данные
    data_cards=[
        ('Обязательно', ['История продаж по дням', 'Минимум 4 месяца', 'Норма: 1 год, идеально: 2 года', 'Один продукт = один столбец']),
        ('Повышает точность', ['Календарь акций: прошлое + план', 'Праздники (календарь известен)', 'Цены по дням', 'Посещаемость (только прошлое)']),
        ('Бонусы', ['Прогноз погоды на будущее', 'Маркетинговые активности', 'Производственные ограничения']),
    ],
    data_panel='Правило будущего',
    data_lines=[
        'Будущее передаётся только для заранее известных фактов: акции, праздники, цены по прайсу.',
        '!Выдуманные будущие числа портят прогноз. Нет плана промо — столбец заполняется нулями, модель честно их не знает.',
        'Формат: CSV/Excel, одна строка = один день, числа с точкой, без строк «итого». Полный чек-лист: docs/data-guide.md.',
    ],

    # 05 · Роли и путь
    wf_steps=[('Данные', 'Продажи выгружают историю,\nмаркетолог даёт план промо'),
              ('Сборка', 'Агент-сборщик: чек-лист,\nсклейка CSV, preflight системы'),
              ('Прогноз', 'Контур A: интервалы 60/80%,\nаномалии, метрики на holdout'),
              ('Решение', 'Ревьюер допускает →\nдиректор планирует по цифрам')],
    roles_panel='Разделение ответственности',
    roles_lines=[
        'Люди: полнота и правдивость данных, план акций, финальное решение по производству.',
        '!Агенты: формат и склейка, запуск модели, ревью качества и лицензии. Агент не передаёт сырые цифры — сначала holdout-метрики.',
        'Каденция: история раз в неделю, прогноз еженедельно, сравнение моделей и проверка лицензии 3.0 — раз в месяц.',
    ],

    # 06 · Архитектура
    arch_items=[
        ('Контур A · ПРОД — TimesFM 2.5, Apache-2.0', 'campaign_forecast.py: точечный прогноз + 60/80% интервалы, XReg-ковариаты (промо, праздники, цены), аномалии, holdout-метрики'),
        ('Контур B · RESEARCH — TimesFM 3.0, non-commercial', 'research_bench.py: нативная мультисерийность, past-future ковариаты за один проход (~6x быстрее), 9 квантилей'),
        ('Лицензионная граница', 'Прогноз 3.0 не попадает в производственные решения — только внутреннее сравнение моделей. Пути к коммерческому 3.0: Apache-релиз, BigQuery AI.FORECAST, прямая лицензия'),
    ],

    # 07 · Лицензия
    lic_cards=[
        ('Разрешено — 3.0', ['Бенчмарк против 2.5', 'Внутренние отчёты о моделях', 'Outputs не являются Derivative', 'Исследования и оценка']),
        ('Запрещено — 3.0', ['План производства по цифрам 3.0', 'Продажа прогноза как услуги', 'Дистилляция в коммерческую модель', 'Дистрибуция весов']),
        ('Коммерческий путь', ['2.5 — Apache-2.0 (уже прод)', 'BigQuery AI.FORECAST — API', 'Прямая лицензия у Google', 'Мониторинг: раз в месяц']),
    ],
    lic_panel='Грань, которую держим в коде',
    lic_lines=[
        'research_bench.py — отдельный вход, не импортируется прод-пайплайном, предупреждение при каждом запуске.',
        '!Грань: цифра 3.0 попала в производственный план = нарушение. Сравнение 3.0 vs 2.5 во внутреннем отчёте = можно.',
    ],

    # 08 · Бенчмарк
    bench_rows=[
        ('MAE, весь горизонт (T=400)', '0.120', '0.107', '0.128', '0.105'),
        ('MAE, обычные дни', '0.129', '0.114', '0.134', '0.108'),
        ('MAE, промо-дни (T=200, короткая история)', '0.412', '0.381', '0.359', '0.142'),
        ('Время, ковариатный прогноз', '—', '~1.3 с', '~0.2 с', '~0.23 с (6x)'),
    ],
    bench_head=('Метрика', '2.5 база', '2.5 XReg', '3.0 база', '3.0 ковар'),
    bench_note='Живой прогон 04.09.2026 на этой машине (Ryzen AI 9 H 365, 20 CPU): синтетика с сезонностью и промо +20%. Реальный ряд из репо (106 точек): 2.5+XReg MAE 31.4 vs 3.0+ковариат 40.5 — на короткой истории 2.5 выигрывает.',

    # 09 · Выход прогноза
    out_head=('Дата', 'Forecast', 'Cautious', 'Optimistic', 'Промо'),
    out_rows=[
        ('2026-09-05', '196.8', '184.2', '209.5', 'акция'),
        ('2026-09-06', '204.1', '191.0', '217.3', 'акция'),
        ('2026-09-07', '172.5', '160.8', '184.1', '—'),
        ('2026-09-08', '141.2', '129.6', '152.7', '—'),
    ],
    out_panel_head='Как читать и что делать',
    out_lines=[
        '!forecast — план производства/закупки по умолчанию. cautious (10-й перцентиль) — нижняя граница: скоропорт ближе к ней. optimistic (90-й) — запас для заморозки и долгих сроков.',
        'В промо-дни (5–6 сентября) модель подняла прогноз заранее — календарь акций из плана сработал.',
        'Метрики на holdout: MAE/RMSE/MAPE + покрытие 80% интервала (цель ~80%). MAPE > 15% — сигнал усилить данные: длиннее история, добавить ковариаты.',
    ],

    # 10 · Каденция
    cad_rows=[
        ('Еженедельно', 'Выгрузка истории за прошедшую неделю', 'Продажи → агент-сборщик'),
        ('Еженедельно', 'Прогноз на 2 недели вперёд', 'Агент-прогнозист (контур A)'),
        ('Раз в месяц', 'Сравнение 2.5 vs 3.0 на свежих данных', 'Агент-исследователь (контур B)'),
        ('Раз в месяц', 'Проверка: Apache-релиз 3.0? BigQuery?', 'Агент-исследователь'),
        ('По событию', 'Внеплановый прогноз под акцию/запуск', 'По заявке директора'),
    ],

    # 11 · Roadmap
    roadmap_items=[
        ('Ближайшее', ['Прогноз на реальных данных мясопродуктов (календарь акций + праздники + цены)', 'Подключение 1С-выгрузки к формату history.csv']),
        ('Средний горизонт', ['BigQuery AI.FORECAST, когда 3.0 появится — облачный контур для масштаба', 'Мультивариантный режим 3.0 в прод-пайплайне после Apache-релиза']),
        ('Экосистема Vector', ['Связка с vector-work (роли) и vector-marketing (кампании)', 'Единый прогнозный слой для всех продуктов Osmosy']),
    ],
    fact_big='2 контура', fact_small='прод на Apache-2.0 и research — лицензионная чистота встроена в архитектуру',
    final_msg='Прогноз — инструмент. Решение — человек.',
    final_sub='github.com/Osmosy/vector-prediction',
)

from engine import kicker, title_block, footer, bg_fill, card, chip, wide_panel, \
    add_text, add_rect, add_bullets
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

ASSETS = '/tmp/vl_assets'

def _rgb(h): return RGBColor.from_string(h)


def _table(slide, th, x, y, col_ws, head, rows, row_h=0.46, head_h=0.42,
           hl_col=None, size=11):
    """Таблица: заголовок + строки; hl_col — индекс колонки, подсвеченной акцентом."""
    xs = [x]
    for w in col_ws[:-1]:
        xs.append(xs[-1] + w)
    total_w = sum(col_ws)
    # header
    add_rect(slide, x, y, total_w, head_h, fill=th['surface'], line=th['card_line'],
             line_w=1.0, radius=th['radius'])
    for j, htxt in enumerate(head):
        add_text(slide, xs[j] + 0.14, y + 0.10, col_ws[j] - 0.2, 0.3, htxt, size=size,
                 bold=True, font=th['f_mono'], color=th['kicker_accent'])
    yy = y + head_h + 0.06
    for i, row in enumerate(rows):
        add_rect(slide, x, yy, total_w, row_h, fill=th['card'], line=th['card_line'],
                 line_w=0.75, radius=th['radius'])
        for j, cell in enumerate(row):
            emph = (j == hl_col)
            add_text(slide, xs[j] + 0.14, yy + 0.11, col_ws[j] - 0.2, 0.3, str(cell),
                     size=size, bold=emph, font=th['f_body'] if j == 0 else th['f_mono'],
                     color=th['accent'] if emph else th['text'])
        yy += row_h + 0.06
    return yy


def sl_title(slide, th, D):
    if th.get('art'):
        from engine import full_art, half_art
        if th['art'][0] == 'full':
            full_art(slide, th, th['art'][1], overlay_pct=26)
        else:
            half_art(slide, th, th['art'][1])
    ta = th['title_align']
    if ta == 'center' and not (th.get('art') and th['art'][0] == 'half'):
        add_rect(slide, (13.333-4.6)/2, 1.52, 4.6, 0.42, fill=th['surface'],
                 line=th['card_line'], radius=0.21, alpha=70)
        add_text(slide, (13.333-4.6)/2, 1.60, 4.6, 0.3, D['kicker'], size=10,
                 bold=True, font=th['f_mono'], color=th['text'],
                 align=PP_ALIGN.CENTER, spacing=140)
        add_text(slide, 1.67, 2.02, 10.0, 1.2, D['title'], size=60, bold=True,
                 font=th['f_title'], color=th['accent'], align=PP_ALIGN.CENTER)
        add_text(slide, 2.67, 3.28, 8.0, 0.7, D['subtitle'], size=15.5,
                 font=th['f_body'], color=th['text'], align=PP_ALIGN.CENTER)
        for i, (big, small) in enumerate(D['chips']):
            chip(slide, th, 3.47 + i*2.25, 4.12, 1.95, big, small)
        add_rect(slide, 4.87, 5.62, 3.6, 0.52, fill=th['accent'], radius=0.26)
        add_text(slide, 4.87, 5.74, 3.6, 0.3, D['github'], size=11.5, bold=True,
                 font=th['f_mono'], color='FFFFFF', align=PP_ALIGN.CENTER)
        add_text(slide, 4.17, 6.90, 5.0, 0.26, D['footer_tag'], size=9.5,
                 font=th['f_mono'], color=th['muted'], align=PP_ALIGN.CENTER)
    else:
        add_rect(slide, 0.62, 1.30, 0.05, 2.2, fill=th['accent'])
        add_text(slide, 0.92, 1.30, 6.5, 0.3, D['kicker'], size=10.5, bold=True,
                 font=th['f_mono'], color=th['kicker_accent'], spacing=160)
        add_text(slide, 0.88, 1.66, 6.2, 1.15, D['title'], size=54, bold=True,
                 font=th['f_title'], color=th['accent'])
        add_text(slide, 0.92, 2.86, 5.9, 0.7, D['subtitle'], size=15,
                 font=th['f_body'], color=th['text'])
        for i, (big, small) in enumerate(D['chips']):
            x = 0.92 + i*1.95
            add_text(slide, x, 3.85, 1.8, 0.5, big, size=26, bold=True,
                     font=th['f_title'], color=th['accent'])
            add_text(slide, x, 4.36, 1.8, 0.3, small, size=9.5, font=th['f_mono'],
                     color=th['muted'], spacing=100)
        add_text(slide, 0.92, 5.15, 6.0, 0.3, D['github'], size=12.5, bold=True,
                 font=th['f_mono'], color=th['accent'])
        add_text(slide, 0.92, 6.0, 6.0, 0.26, D['footer_tag'], size=9.5,
                 font=th['f_mono'], color=th['muted'])
    slide.shapes.add_picture(f'{ASSETS}/vector_ray_t.png',
                             Inches(12.20), Inches(0.30), width=Inches(0.9))


def _three_cards(slide, th, D, cards_key, head_key, lead_key, panel_key, panel_lines_key, num):
    bg_fill(slide, th)
    kicker(slide, th, f'{num:02d} · ')
    return None  # не используется


def _content_slide(slide, th, D, num, head, lead, cards, panel_head, panel_lines):
    bg_fill(slide, th)
    kicker(slide, th, f'{num:02d} · ')
    title_block(slide, th, head)
    if lead:
        add_text(slide, 0.62, 1.75, 12.1, 0.5, lead, size=13,
                 font=th['f_body'], color=th['muted'])
    cw, gap, m = 3.95, 0.25, 0.62
    for i, (ch, lines) in enumerate(cards):
        card(slide, th, m + i*(cw+gap), 2.42 if lead else 2.15, cw, 2.9 if lead else 3.2,
             ch, lines)
    if panel_head:
        wide_panel(slide, th, 0.62, 5.62 if lead else 5.62, 12.1,
                   1.75 if panel_lines and len(panel_lines) > 2 else 1.45,
                   panel_head, panel_lines)


def sl_intro(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '01 · Введение')
    title_block(slide, th, 'Что такое Vector Prediction')
    add_text(slide, 0.62, 1.75, 12.1, 0.5, D['intro_lead'], size=13,
             font=th['f_body'], color=th['muted'])
    cw, gap, m = 3.95, 0.25, 0.62
    for i, (head, lines) in enumerate(D['intro_cards']):
        card(slide, th, m + i*(cw+gap), 2.55, cw, 3.4, head, lines)
    footer(slide, th, 2)


def sl_problem(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '02 · Зачем')
    title_block(slide, th, 'Прогноз, которого обычно нет')
    cw, gap, m = 3.95, 0.25, 0.62
    for i, (head, lines) in enumerate(D['prob_cards']):
        card(slide, th, m + i*(cw+gap), 2.15, cw, 3.1, head, lines)
    wide_panel(slide, th, 0.62, 5.55, 12.1, 1.85, D['prob_panel'], D['prob_lines'])
    footer(slide, th, 3)


def sl_solution(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '03 · Решение')
    title_block(slide, th, 'Foundation model вместо ручных таблиц')
    cw, gap, m = 3.95, 0.25, 0.62
    for i, (head, lines) in enumerate(D['sol_cards']):
        card(slide, th, m + i*(cw+gap), 2.15, cw, 3.1, head, lines)
    wide_panel(slide, th, 0.62, 5.55, 12.1, 1.85, D['sol_panel'], D['sol_lines'])
    footer(slide, th, 4)


def sl_data(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '04 · Данные')
    title_block(slide, th, 'Что подготовить: три уровня')
    cw, gap, m = 3.95, 0.25, 0.62
    for i, (head, lines) in enumerate(D['data_cards']):
        card(slide, th, m + i*(cw+gap), 2.15, cw, 3.1, head, lines)
    wide_panel(slide, th, 0.62, 5.55, 12.1, 1.85, D['data_panel'], D['data_lines'])
    footer(slide, th, 5)


def sl_roles(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '05 · Кто что делает')
    title_block(slide, th, 'Люди решают — агенты считают')
    xs = [0.55, 3.80, 7.05, 10.30]
    for i, (h, l) in enumerate(D['wf_steps']):
        card(slide, th, xs[i], 2.10, 2.92, 2.0, h, l.split('\n'), num=i + 1)
        if i < 3:
            add_rect(slide, xs[i] + 2.97, 3.05, 0.28, 0.035, fill=th['accent'])
    wide_panel(slide, th, 0.83, 4.70, 11.67, 2.1, D['roles_panel'], D['roles_lines'])
    footer(slide, th, 6)


def sl_arch(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '06 · Архитектура')
    title_block(slide, th, 'Двухконтурная схема')
    y = 2.2
    hs = [1.55, 1.55, 1.85]
    for i, (head, lines) in enumerate(D['arch_items']):
        wide_panel(slide, th, 0.62, y, 12.1, hs[i], head, [lines])
        y += hs[i] + 0.18
    footer(slide, th, 7)


def sl_license(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '07 · Лицензия TimesFM 3.0')
    title_block(slide, th, 'Non-commercial: что можно и нельзя')
    cw, gap, m = 3.95, 0.25, 0.62
    for i, (head, lines) in enumerate(D['lic_cards']):
        card(slide, th, m + i*(cw+gap), 2.15, cw, 3.1, head, lines)
    wide_panel(slide, th, 0.62, 5.55, 12.1, 1.7, D['lic_panel'], D['lic_lines'])
    footer(slide, th, 8)


def sl_bench(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '08 · Бенчмарк 2.5 vs 3.0')
    title_block(slide, th, 'Точность и скорость на живом прогоне')
    _table(slide, th, 0.62, 2.05, [3.9, 1.55, 1.55, 1.55, 1.55],
           D['bench_head'], D['bench_rows'], row_h=0.52, hl_col=4)
    wide_panel(slide, th, 0.62, 5.30, 12.1, 1.95, 'Как читать', [
        '!Вывод: точность сопоставима (лучший — 3.0 с ковариатом), скорость ковариатного прогноза у 3.0 в ~6 раз выше.',
        'Для прод-календаря акций 2.5+XReg достаточен: на выраженном промо-паттерне ошибка промо-дня близка к нулю.',
        D['bench_note'],
    ])
    footer(slide, th, 9)


def sl_output(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '09 · Результат')
    title_block(slide, th, 'Что получает директор')
    _table(slide, th, 0.62, 2.05, [2.2, 1.9, 1.9, 1.9, 1.6],
           D['out_head'], D['out_rows'], row_h=0.52, hl_col=1)
    wide_panel(slide, th, 0.62, 4.85, 12.1, 2.35, D['out_panel_head'], D['out_lines'])
    footer(slide, th, 10)


def sl_cadence(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '10 · Каденция')
    title_block(slide, th, 'Ритм работы прогнозного контура')
    _table(slide, th, 0.62, 2.05, [2.2, 5.6, 4.3],
           ('Цикл', 'Что происходит', 'Кто'), D['cad_rows'], row_h=0.62)
    wide_panel(slide, th, 0.62, 6.05, 12.1, 1.2, 'Правило контура B', [
        'research_bench.py не встраивается в прод-пайплайн: отдельный вход, отдельный выход.',
    ])
    footer(slide, th, 11)


def sl_roadmap(slide, th, D):
    bg_fill(slide, th)
    kicker(slide, th, '11 · Развитие')
    title_block(slide, th, 'Roadmap')
    ys = [1.95, 3.15, 4.75]
    hts = [1.0, 1.4, 1.1]
    for (h, lines), y, hh in zip(D['roadmap_items'], ys, hts):
        wide_panel(slide, th, 0.62, y, 12.1, hh, h, lines)
    footer(slide, th, 12)


def sl_final(slide, th, D):
    bg_fill(slide, th)
    add_text(slide, 1.67, 2.0, 10.0, 1.0, D['fact_big'], size=58, bold=True,
             font=th['f_title'], color=th['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, 1.17, 3.25, 11.0, 0.5, D['fact_small'], size=14,
             font=th['f_body'], color=th['text'], align=PP_ALIGN.CENTER)
    add_text(slide, 2.67, 4.35, 8.0, 0.6, D['final_msg'], size=22, bold=True,
             font=th['f_title'], color=th['text'], align=PP_ALIGN.CENTER)
    add_rect(slide, 5.17, 5.25, 3.0, 0.035, fill=th['accent'])
    add_text(slide, 2.67, 5.65, 8.0, 0.3, D['final_sub'], size=12.5,
             font=th['f_mono'], color=th['accent'], align=PP_ALIGN.CENTER)
    add_text(slide, 2.67, 6.85, 8.0, 0.24,
             'Контур A: Apache-2.0 (прод) · Контур B: non-commercial (research only) · гайды и лицензия — в репо',
             size=9.5, font=th['f_body'], color=th['muted'], align=PP_ALIGN.CENTER)
    slide.shapes.add_picture(f'{ASSETS}/vector_ray_t.png',
                             Inches(12.20), Inches(0.30), width=Inches(0.9))


SLIDES = [sl_title, sl_intro, sl_problem, sl_solution, sl_data, sl_roles,
          sl_arch, sl_license, sl_bench, sl_output, sl_cadence, sl_roadmap, sl_final]